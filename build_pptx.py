#!/usr/bin/env python3
"""Build 'How AI is Changing High School Education' — 6-slide .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──
OUT_DIR = r"C:\Users\LENOVO\projects\ai-education-pptx"
OUT_FILE = os.path.join(OUT_DIR, "AI-Education-Claude.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

# ── Palette ──
DEEP_PURPLE  = RGBColor(0x5B, 0x21, 0xB6)
INDIGO       = RGBColor(0x43, 0x38, 0xCA)
ACCENT_GOLD  = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GOLD2 = RGBColor(0xF7, 0xAB, 0x2E)
DARK_BG      = RGBColor(0x1E, 0x0A, 0x3E)  # near-black-purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG     = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x0F, 0x17, 0x2A)
MED_TEXT      = RGBColor(0x47, 0x55, 0x69)
SOFT_BORDER   = RGBColor(0xE2, 0xE8, 0xF0)
BENEFIT_GREEN = RGBColor(0x10, 0xB9, 0x81)
RISK_RED      = RGBColor(0xEF, 0x44, 0x44)

# ── Helpers ──
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_name="Calibri",
                font_size=Pt(18), bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = Pt(2)
    if line_spacing:
        p.line_spacing = line_spacing
    return tf

def add_rich_textbox(slide, left, top, width, height, paragraphs_data):
    """paragraphs_data: list of dicts with text, font_name, font_size, bold, color, alignment"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pdict in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pdict.get("text", "")
        p.font.name = pdict.get("font_name", "Calibri")
        p.font.size = pdict.get("font_size", Pt(14))
        p.font.bold = pdict.get("bold", False)
        p.font.color.rgb = pdict.get("color", DARK_TEXT)
        p.alignment = pdict.get("alignment", PP_ALIGN.LEFT)
        p.space_after = pdict.get("space_after", Pt(4))
        if "line_spacing" in pdict:
            p.line_spacing = pdict["line_spacing"]
    return tf

def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=SOFT_BORDER, border_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    shape.shadow.inherit = False
    return shape

def add_accent_bar(slide, left, top, width, height, color=ACCENT_GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ── Presentation ──
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Slide Layout (blank) ──
blank_layout = prs.slide_layouts[6]  # blank

team_names = ["Alex Chen", "Jordan Rivera", "Sam Park", "Taylor Brooks", "Morgan Liu"]

# ═══════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE (dark)
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, DARK_BG)

# Decorative accent shapes
add_accent_bar(slide1, 0, 0, 13.333, 0.08, ACCENT_GOLD)
add_accent_bar(slide1, 0, 2.3, 0.12, 2.5, ACCENT_GOLD)

# Title
add_textbox(slide1, 1.2, 2.4, 11.0, 1.0,
            "How AI is Changing", "Arial Black", Pt(44), True, WHITE, PP_ALIGN.LEFT)
add_textbox(slide1, 1.2, 3.15, 11.0, 0.8,
            "High School Education", "Arial Black", Pt(44), True, ACCENT_GOLD, PP_ALIGN.LEFT)

# Subtitle line
add_textbox(slide1, 1.2, 4.2, 10.0, 0.5,
            "An exploration of artificial intelligence in today's classrooms",
            "Calibri", Pt(18), False, RGBColor(0xA5, 0xB4, 0xFC), PP_ALIGN.LEFT)

# Team names
team_text = "  •  ".join(team_names)
add_textbox(slide1, 1.2, 5.3, 11.0, 0.5,
            team_text,
            "Calibri", Pt(16), False, RGBColor(0x94, 0xA3, 0xB8), PP_ALIGN.LEFT)

# Bottom accent
add_accent_bar(slide1, 0, 7.42, 13.333, 0.08, ACCENT_GOLD)

# ═══════════════════════════════════════════
# SLIDE 2 — WHAT IS AI? (light)
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, LIGHT_BG)

# Top bar
add_accent_bar(slide2, 0, 0, 13.333, 0.06, DEEP_PURPLE)

# Section label
add_textbox(slide2, 0.8, 0.35, 3.0, 0.45, "WHAT IS AI?",
            "Arial Black", Pt(14), True, DEEP_PURPLE, PP_ALIGN.LEFT)

# Title
add_textbox(slide2, 0.8, 0.85, 11.5, 0.7,
            "Understanding Artificial Intelligence",
            "Arial Black", Pt(32), True, DARK_TEXT, PP_ALIGN.LEFT)

# Gold underline
add_accent_bar(slide2, 0.8, 1.6, 2.0, 0.06, ACCENT_GOLD)

# Cards row
card_data = [
    ("🤖", "Machines that Learn",
     "AI refers to computer systems\nthat can perform tasks normally\nrequiring human intelligence —\nlearning, problem-solving,\nand pattern recognition."),
    ("🧠", "How It Works",
     "AI uses algorithms and\nlarge amounts of data to\nrecognize patterns, make\npredictions, and improve\nover time through training."),
    ("💡", "Everyday Examples",
     "You already use AI daily:\nvoice assistants (Siri/Alexa),\nrecommendation systems (Netflix,\nSpotify), navigation apps,\nand smart replies in messaging."),
]

for idx, (emoji, title, body) in enumerate(card_data):
    left = 0.8 + idx * 4.1
    add_card(slide2, left, 2.1, 3.8, 3.4)

    # Top color bar on card
    add_accent_bar(slide2, left, 2.1, 3.8, 0.06, INDIGO if idx != 1 else ACCENT_GOLD)

    add_textbox(slide2, left + 0.3, 2.35, 3.2, 0.5, emoji, "Segoe UI Emoji", Pt(28), False, DARK_TEXT)

    add_textbox(slide2, left + 0.3, 2.85, 3.2, 0.4, title,
                "Arial Black", Pt(16), True, DEEP_PURPLE)

    add_textbox(slide2, left + 0.3, 3.35, 3.2, 1.9, body,
                "Calibri", Pt(13), False, MED_TEXT, line_spacing=1.3)

# Fun fact at bottom
add_card(slide2, 0.8, 5.8, 11.7, 1.2, RGBColor(0xED, 0xE9, 0xFE), RGBColor(0xC4, 0xB5, 0xFD), Pt(1.5))
add_textbox(slide2, 1.1, 5.95, 11.0, 0.9,
            '⚡  Did you know?  The term "Artificial Intelligence" was coined in 1956 — but it\'s only in the last decade that AI has become powerful enough to transform how we learn, work, and live.',
            "Calibri", Pt(14), False, DEEP_PURPLE, PP_ALIGN.LEFT, line_spacing=1.3)

# ═══════════════════════════════════════════
# SLIDE 3 — AI IN LEARNING (light)
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, LIGHT_BG)

add_accent_bar(slide3, 0, 0, 13.333, 0.06, DEEP_PURPLE)
add_textbox(slide3, 0.8, 0.35, 3.0, 0.45, "APPLICATIONS",
            "Arial Black", Pt(14), True, DEEP_PURPLE)
add_textbox(slide3, 0.8, 0.85, 11.5, 0.7,
            "AI in the Learning Process",
            "Arial Black", Pt(32), True, DARK_TEXT)
add_accent_bar(slide3, 0.8, 1.6, 2.0, 0.06, ACCENT_GOLD)

# Three feature blocks — horizontal cards with accent dots
features = [
    ("🎯", "Personalized Tutoring",
     "AI adapts to each student's pace\nand learning style. Tools like\nKhan Academy's Khanmigo offer\none-on-one tutoring tailored to\nwhere you're struggling most.",
     DEEP_PURPLE),
    ("⚡", "Instant Feedback",
     "Unlike waiting days for a\nteacher to grade papers, AI\nprovides immediate feedback on\nwriting, math problems, and\ncoding exercises — helping you\nimprove in real time.",
     INDIGO),
    ("🌍", "Language Learning",
     "Apps like Duolingo use AI to\ncreate adaptive lessons, provide\npronunciation feedback, and\nsimulate real conversations —\nmaking language learning more\nengaging and effective.",
     ACCENT_GOLD),
]

for idx, (emoji, title, body, color) in enumerate(features):
    left = 0.8 + idx * 4.1
    add_card(slide3, left, 2.1, 3.8, 3.6)

    # Accent dot
    add_circle(slide3, left + 0.3, 2.35, 0.45, color)

    add_textbox(slide3, left + 1.0, 2.35, 2.5, 0.4, emoji, "Segoe UI Emoji", Pt(22), False, DARK_TEXT)

    add_textbox(slide3, left + 0.3, 2.9, 3.2, 0.4, title,
                "Arial Black", Pt(17), True, color)

    add_textbox(slide3, left + 0.3, 3.4, 3.2, 2.0, body,
                "Calibri", Pt(12.5), False, MED_TEXT, line_spacing=1.3)

# Bottom impact stat
add_card(slide3, 0.8, 6.05, 11.7, 1.0, RGBColor(0xFE, 0xFB, 0xEB), RGBColor(0xFD, 0xE6, 0x8A), Pt(1.5))
add_textbox(slide3, 1.1, 6.15, 11.0, 0.8,
            "📊  Research shows AI tutoring can improve student performance by up to 30% compared to traditional classroom instruction alone.",
            "Calibri", Pt(14), False, RGBColor(0x92, 0x4E, 0x0F), PP_ALIGN.LEFT, line_spacing=1.2)

# ═══════════════════════════════════════════
# SLIDE 4 — AI TOOLS STUDENTS USE (light)
# ═══════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, LIGHT_BG)

add_accent_bar(slide4, 0, 0, 13.333, 0.06, DEEP_PURPLE)
add_textbox(slide4, 0.8, 0.35, 3.0, 0.45, "TOOLS",
            "Arial Black", Pt(14), True, DEEP_PURPLE)
add_textbox(slide4, 0.8, 0.85, 11.5, 0.7,
            "AI Tools Students Are Using Today",
            "Arial Black", Pt(32), True, DARK_TEXT)
add_accent_bar(slide4, 0.8, 1.6, 2.0, 0.06, ACCENT_GOLD)

# Four tool cards in 2x2 grid
tools = [
    ("ChatGPT", "💬", "Conversational AI that helps\nwith research, brainstorming,\nwriting assistance, and\nexplaining complex topics\nin simple terms.", DEEP_PURPLE),
    ("Grammarly", "✍️", "AI-powered writing assistant\nthat checks grammar, tone, and\nclarity — helping students become\nbetter writers across every\nsubject.", INDIGO),
    ("Photomath", "📐", "Scan a math problem and get\nstep-by-step explanations.\nGreat for understanding the\n\"how\" behind the answer,\nnot just the answer itself.", ACCENT_GOLD),
    ("Quizlet AI", "📚", "Creates personalized flashcards,\npractice tests, and study plans\nfrom your notes — using AI\nto make studying smarter\nand more efficient.", RGBColor(0x8B, 0x5C, 0xF6)),
]

# 2x2 grid: first row
for idx, (name, emoji, desc, color) in enumerate(tools[:2]):
    left = 0.8 + idx * 6.1
    add_card(slide4, left, 2.1, 5.8, 2.35)
    add_accent_bar(slide4, left, 2.1, 5.8, 0.06, color)

    add_textbox(slide4, left + 0.3, 2.3, 5.2, 0.35, emoji, "Segoe UI Emoji", Pt(22), False, DARK_TEXT)
    add_textbox(slide4, left + 0.8, 2.3, 4.5, 0.35, name,
                "Arial Black", Pt(18), True, color)
    add_textbox(slide4, left + 0.3, 2.8, 5.2, 1.5, desc,
                "Calibri", Pt(12), False, MED_TEXT, line_spacing=1.25)

# Second row
for idx, (name, emoji, desc, color) in enumerate(tools[2:]):
    left = 0.8 + idx * 6.1
    add_card(slide4, left, 4.7, 5.8, 2.35)
    add_accent_bar(slide4, left, 4.7, 5.8, 0.06, color)

    add_textbox(slide4, left + 0.3, 4.9, 5.2, 0.35, emoji, "Segoe UI Emoji", Pt(22), False, DARK_TEXT)
    add_textbox(slide4, left + 0.8, 4.9, 4.5, 0.35, name,
                "Arial Black", Pt(18), True, color)
    add_textbox(slide4, left + 0.3, 5.4, 5.2, 1.5, desc,
                "Calibri", Pt(12), False, MED_TEXT, line_spacing=1.25)

# ═══════════════════════════════════════════
# SLIDE 5 — BENEFITS & CONCERNS (light)
# ═══════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, LIGHT_BG)

add_accent_bar(slide5, 0, 0, 13.333, 0.06, DEEP_PURPLE)
add_textbox(slide5, 0.8, 0.35, 3.0, 0.45, "ANALYSIS",
            "Arial Black", Pt(14), True, DEEP_PURPLE)
add_textbox(slide5, 0.8, 0.85, 11.5, 0.7,
            "Benefits & Concerns",
            "Arial Black", Pt(32), True, DARK_TEXT)
add_accent_bar(slide5, 0.8, 1.6, 2.0, 0.06, ACCENT_GOLD)

# ── LEFT COLUMN: Benefits ──
add_card(slide5, 0.8, 2.1, 5.8, 4.8)
add_accent_bar(slide5, 0.8, 2.1, 5.8, 0.06, BENEFIT_GREEN)

# Benefits header
add_card(slide5, 1.1, 2.3, 2.4, 0.5, RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0xA7, 0xF3, 0xD0), Pt(1))
add_textbox(slide5, 1.3, 2.35, 2.0, 0.4, "✅  BENEFITS",
            "Arial Black", Pt(14), True, BENEFIT_GREEN, PP_ALIGN.CENTER)

benefits = [
    ("Personalized Learning", "Adapts to individual student needs, pace, and learning preferences"),
    ("24/7 Availability", "Students can get help anytime — not just during school hours"),
    ("Engaging Content", "Interactive lessons, gamification, and multimedia make learning fun"),
    ("Teacher Support", "Automates grading and admin tasks, freeing teachers to teach"),
    ("Accessibility", "Text-to-speech, translation, and assistive tools help all learners"),
    ("Data-Driven Insights", "Identifies learning gaps so teachers can target support"),
]

for i, (title, desc) in enumerate(benefits):
    y = 2.95 + i * 0.62
    add_circle(slide5, 1.25, y + 0.02, 0.18, BENEFIT_GREEN)
    add_textbox(slide5, 1.6, y, 4.8, 0.25, title,
                "Arial Black", Pt(12), True, DARK_TEXT)
    add_textbox(slide5, 1.6, y + 0.25, 4.8, 0.3, desc,
                "Calibri", Pt(10), False, MED_TEXT)

# ── RIGHT COLUMN: Concerns ──
add_card(slide5, 6.9, 2.1, 5.8, 4.8)
add_accent_bar(slide5, 6.9, 2.1, 5.8, 0.06, RISK_RED)

add_card(slide5, 7.2, 2.3, 2.4, 0.5, RGBColor(0xFE, 0xE2, 0xE2), RGBColor(0xFE, 0xC8, 0xCA), Pt(1))
add_textbox(slide5, 7.4, 2.35, 2.0, 0.4, "⚠️  CONCERNS",
            "Arial Black", Pt(14), True, RISK_RED, PP_ALIGN.CENTER)

concerns = [
    ("Academic Integrity", "Easy access to AI raises plagiarism and cheating risks"),
    ("Over-Reliance", "Students may skip learning fundamentals by leaning on AI"),
    ("Privacy Issues", "AI tools collect student data — who owns it and how is it used?"),
    ("Equity Gap", "Not all students have equal access to devices and internet"),
    ("Misinformation", "AI can produce confident-sounding but incorrect information"),
    ("Social Isolation", "Screen-heavy learning may reduce peer interaction and collaboration"),
]

for i, (title, desc) in enumerate(concerns):
    y = 2.95 + i * 0.62
    add_circle(slide5, 7.35, y + 0.02, 0.18, RISK_RED)
    add_textbox(slide5, 7.7, y, 4.8, 0.25, title,
                "Arial Black", Pt(12), True, DARK_TEXT)
    add_textbox(slide5, 7.7, y + 0.25, 4.8, 0.3, desc,
                "Calibri", Pt(10), False, MED_TEXT)

# ═══════════════════════════════════════════
# SLIDE 6 — KEY TAKEAWAYS + Q&A (dark)
# ═══════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, DARK_BG)

add_accent_bar(slide6, 0, 0, 13.333, 0.08, ACCENT_GOLD)

add_textbox(slide6, 0.8, 0.5, 3.0, 0.45, "CONCLUSION",
            "Arial Black", Pt(14), True, RGBColor(0xA5, 0xB4, 0xFC))
add_textbox(slide6, 0.8, 1.05, 11.5, 0.8,
            "Key Takeaways",
            "Arial Black", Pt(36), True, WHITE)
add_accent_bar(slide6, 0.8, 1.9, 2.0, 0.06, ACCENT_GOLD)

takeaways = [
    ("1", "AI is already here — and it's transforming how high school\nstudents learn, study, and prepare for the future."),
    ("2", "From personalized tutoring to instant feedback, AI tools are\nmaking education more effective and accessible."),
    ("3", "The key is balance: leveraging AI's benefits while staying\nmindful of its risks and limitations."),
    ("4", "Digital literacy and AI awareness are essential skills for\nsuccess in the 21st-century world."),
]

for i, (num, text) in enumerate(takeaways):
    y = 2.3 + i * 0.95
    # Number circle
    add_circle(slide6, 1.0, y + 0.05, 0.55, ACCENT_GOLD)
    add_textbox(slide6, 1.0, y + 0.1, 0.55, 0.45, num,
                "Arial Black", Pt(20), True, DARK_BG, PP_ALIGN.CENTER)
    # Text
    add_textbox(slide6, 1.8, y + 0.02, 10.5, 0.8, text,
                "Calibri", Pt(16), False, RGBColor(0xE2, 0xE8, 0xF0), line_spacing=1.3)

# Q&A section
add_card(slide6, 3.0, 6.0, 7.3, 1.1, DEEP_PURPLE, ACCENT_GOLD, Pt(2))
add_textbox(slide6, 3.3, 6.15, 6.7, 0.45, "🙋  Questions & Discussion",
            "Arial Black", Pt(22), True, ACCENT_GOLD, PP_ALIGN.CENTER)
add_textbox(slide6, 3.3, 6.6, 6.7, 0.35,
            "Thank you for listening!  —  Alex, Jordan, Sam, Taylor & Morgan",
            "Calibri", Pt(14), False, RGBColor(0xC7, 0xD2, 0xFE), PP_ALIGN.CENTER)

add_accent_bar(slide6, 0, 7.42, 13.333, 0.08, ACCENT_GOLD)

# ── Save ──
prs.save(OUT_FILE)
print(f"✅ Presentation saved to: {OUT_FILE}")
